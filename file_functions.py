import os
import sys
import csv
import pandas as pd
from pptx.util import Pt
from pptx.util import Inches
from pptx import Presentation
from Errors import throw
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from openpyxl import Workbook

class FileFunctions():
    """
        This class will read and check
        files for further use.
    """

    def __init__(self, file):
        self.file = file

    def file_exists(self):
        """
            This function checks if a file exists and
            if the program can read the content.
        """
        if os.path.exists(self.file):
            print('\033[92mFile specified exists !\033[0m\n')
        else:
            throw('File specified not found, exiting.', 'high')

    def read_file(self):
        """
            This function will return the content of
            a normal file.
        """
        try:
            file = open(self.file, 'r', encoding='latin-1')
            text = file.read()
            file.close()
        except OSError:
            throw("Couldn't read file, exiting.", "high")       
        return text
    
    def read_log_file(self):
        """This function returns the content of a hardening log file
        """
        log_content = pd.read_table(self.file, encoding='latin-1')
        log_content.columns = ["LOG"]
        return log_content

    def read_csv_file(self):
        """
            This function will return a dataframe (pandas)
            containing the whole data of a CSV file.
        """
        try:
            df = pd.read_csv(self.file, encoding='latin1')
        except pd.errors.ParserError:
            throw('An error occured while reading csv, please make sure it\'s using a good format', 'error')
        df = df.fillna('')
        df = df.astype(str)
        return df

    def read_xlsx_tracefile(self):
        """
            This function will return Excel sheets from trace file
        """
        df_all_policies = pd.read_excel(self.file, "All-Policies").fillna('')
        df_contexts = pd.read_excel(self.file, "Contexts").fillna('')

        return df_all_policies, df_contexts

    def read_xlsx_contexts_sheet(self):
        """
            This function will retreive information in xlsx
        """
        df_contexts = pd.read_excel(self.file, "Contexts", header=1).fillna('')

        return df_contexts

    def convert_csv_2_excel(self):
        """
            This function will transform a CSV file
            into an Excel file, using pandas.
        """
        df = pd.read_csv(self.file)
        df = df.fillna('')
        output_excel = ''
        output_excel_args = ['-o', '--output']
        for output_excel_arg in output_excel_args:
            for arg in sys.argv:
                if output_excel_arg == arg:
                    output_excel = sys.argv[sys.argv.index(arg)+1]
        if output_excel == '':
            output_excel = input("\nWhat's the name of the Excel output file ? : ")
        df.to_excel(output_excel, index=False)

    def convert_excel_2_csv(self):
        """
            This function will transform an Excel file
            into a CSV file, using pandas.
        """
        df = pd.read_excel(self.file)
        df = df.fillna('')
        output_csv = ''
        output_csv_args = ['-o', '--output']
        for output_csv_arg in output_csv_args:
            for arg in sys.argv:
                if output_csv_arg == arg:
                    output_csv = sys.argv[sys.argv.index(arg)+1]
        if output_csv == '':
            output_csv = input("\nWhat's the name of the CSV output file ? : ")
        df.to_csv(output_csv, index=False)

    def iter_cells(self, table):
        """
            This function will itterate over a PowerPoint
            table to transform the content.
        """
        for row in table.rows:
            for cell in row.cells:
                yield cell

    def create_powerpoint(self, hardening_dataframe: pd.DataFrame, contexts: list
    , contexts_columns: list, powerpoint_filepath: str):
        """
            This function will transform a CSV file
            into PowerPoint Slides.
        """
        # Creating presentation
        prs = Presentation()
        slide_size = (16, 9)
        prs.slide_width, prs.slide_height = Inches(slide_size[0]), Inches(slide_size[1])

        # Creating first slide 
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.width = Inches(16)
        title.height = Inches(1.5)
        title.top = Inches(3)
        subtitle.width = Inches(16)
        subtitle.height = Inches(1.5)
        subtitle.top = Inches(5)

        title.text = "Hardening presentation"
        subtitle.text = "Author"

        # Creating policies slides 
        for _, policy in hardening_dataframe.iterrows():
            tab_slide_layout = prs.slide_layouts[6]
            tab_slide = prs.slides.add_slide(tab_slide_layout)
            shapes = tab_slide.shapes

            # Add Impact, Description and Rationale in notes
            comment_text = ""
            if policy['Impact']:
                comment_text += f"""
Impact :
{policy['Impact']}
"""
            if policy['Description']:
                comment_text += f"""
Description :
{policy['Description']}
"""
            if policy['Rationale']:
                comment_text += f"""
Rationale :
{policy['Rationale']}
"""
            notes_text_frame = tab_slide.notes_slide.notes_text_frame
            notes_text_frame.text = comment_text

            cryptonit_color = RGBColor(57, 100, 0) # dark green

            # Place policy ID
            left = Inches(1.5)
            top = Inches(0.4)
            width = Inches(3.5)
            height = Inches(0.7)
            id_box = shapes.add_textbox(left, top, width, height)
            paragraph = id_box.text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.LEFT
            run = paragraph.add_run()
            run.text = policy['ID']
            run.font.name = "Calibri Light"
            run.font.color.rgb = cryptonit_color
            run.font.size = Pt(36)
            run.font.bold = True
            run.font.underline = True

            # Place policy Severity
            if policy['Severity'] in ['High', 'Medium', 'Low']:
                left = Inches(5)
                top = Inches(0.55)
                width = Inches(1.2)
                height = Inches(0.4)
                severity_rect = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
                severity_rect.shadow.inherit = False
                paragraph = severity_rect.text_frame.paragraphs[0]
                paragraph.alignment = PP_ALIGN.CENTER
                run = paragraph.add_run()
                run.text = policy['Severity']
                run.font.name = "Calibri Light"
                run.font.size = Pt(20)
                severity_rect.fill.solid()
                low_color = RGBColor(0, 112, 192) # blue
                medium_color = RGBColor(239, 169, 6) # orange
                high_color = RGBColor(192, 0, 0) # red

                if policy['Severity'] == "Low":
                    color = low_color
                elif policy['Severity'] == "Medium":
                    color = medium_color
                elif policy['Severity'] == "High":
                    color = high_color

                severity_rect.fill.fore_color.rgb = color
                severity_rect.line.color.rgb = color

            # Place policy Level
            if 'Level' in hardening_dataframe.columns and policy['Level'] in ['NG', 'L1', 'L2']:
                color_NG = RGBColor(146, 145, 157) # grey
                color_L1 = RGBColor(232, 175, 207) # pink
                color_L2 = RGBColor(187, 153, 232) # purple
                left = Inches(6.4)
                top = Inches(0.55)
                width = Inches(0.6)
                height = Inches(0.4)
                level_rect = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
                level_rect.shadow.inherit = False
                paragraph = level_rect.text_frame.paragraphs[0]
                paragraph.alignment = PP_ALIGN.CENTER
                run = paragraph.add_run()
                run.text = policy['Level']
                run.font.name = "Calibri Light"
                run.font.size = Pt(20)
                level_rect.fill.solid()
                if policy['Level'] == 'L1':
                    level_color = color_L1
                elif policy['Level'] == 'L2':
                    level_color = color_L2
                else:
                    level_color = color_NG
                level_rect.fill.fore_color.rgb = level_color
                level_rect.line.color.rgb = level_color

            # Place policy name
            left = Inches(1.5)
            top = Inches(1)
            width = Inches(14.5)
            height = Inches(0.5)
            name_box = shapes.add_textbox(left, top, width, height)
            paragraph = name_box.text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.LEFT
            run = paragraph.add_run()
            run.text = policy['Name']
            run.font.name = "Calibri Light"
            run.font.color.rgb = cryptonit_color
            run.font.size = Pt(30)

            # Add value table
            if len(contexts) == 0:
                cols = 2
            else:
                cols = 5

            rows = 4 + len(contexts)
            left = Inches(2.0)
            top = Inches(1.7)
            width = Inches(13)
            height = Inches(0.12)

            shape = shapes.add_table(rows, cols, left, top, width, height)
            table = shape.table
            table_style= '{EB344D84-9AFB-497E-A393-DC336BA19D2E}'
            tbl =  shape._element.graphic.graphicData.tbl
            tbl[0][-1].text = table_style

            # Change font size
            for cell in self.iter_cells(table):
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(16)

            # set columns widths
            table.columns[0].width = Inches(3)
            table.columns[1].width = Inches(8)
            if len(contexts) > 0:
                table.columns[2].width = Inches(0.75)
                table.columns[3].width = Inches(0.75)
                table.columns[4].width = Inches(0.75)


            # set rows height
            table.rows[0].height = Inches(0.7) # Title
            table.rows[1].height = Inches(0.7) # Possible values
            table.rows[2].height = Inches(0.7) # Default value
            table.rows[3].height = Inches(0.7) # Recommended value
            
            i = 0
            while i < len(contexts):
                table.rows[i + 4].height = Inches(0.7)
                i+=1

            # Default column headings
            table.cell(0, 0).text = 'Variable'
            table.cell(0, 1).text = 'Value'
            table.cell(0, 0).fill.solid()
            table.cell(0, 0).fill.fore_color.rgb = cryptonit_color
            table.cell(0, 1).fill.solid()
            table.cell(0, 1).fill.fore_color.rgb = cryptonit_color

            if len(contexts) > 0:
                table.cell(0, 2).fill.solid()
                table.cell(0, 2).fill.fore_color.rgb = cryptonit_color
                table.cell(0, 3).fill.solid()
                table.cell(0, 3).fill.fore_color.rgb = cryptonit_color
                table.cell(0, 4).fill.solid()
                table.cell(0, 4).fill.fore_color.rgb = cryptonit_color

                table.cell(0, 1).merge(table.cell(0, 4))
                table.cell(1, 1).merge(table.cell(1, 4))
                table.cell(2, 1).merge(table.cell(2, 4))
                #table.cell(3, 2).merge(table.cell(3, 4))

                # create choice result table
                table.cell(3, 2).fill.solid()
                table.cell(3, 2).fill.fore_color.rgb = cryptonit_color
                table.cell(3, 2).text = 'RecVal'
                table.cell(3, 2).text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                table.cell(3, 2).text_frame.paragraphs[0].font.bold = True
                table.cell(3, 2).text_frame.paragraphs[0].font.size = Pt(14)

                table.cell(3, 3).fill.solid()
                table.cell(3, 3).fill.fore_color.rgb = cryptonit_color
                table.cell(3, 3).text = 'Same'
                table.cell(3, 3).text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                table.cell(3, 3).text_frame.paragraphs[0].font.bold = True
                table.cell(3, 3).text_frame.paragraphs[0].font.size = Pt(14)

                table.cell(3, 4).fill.solid()
                table.cell(3, 4).fill.fore_color.rgb = cryptonit_color
                table.cell(3, 4).text = 'Other'
                table.cell(3, 4).text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                table.cell(3, 4).text_frame.paragraphs[0].font.bold = True
                table.cell(3, 4).text_frame.paragraphs[0].font.size = Pt(14)

            # Default 1st column cells
            table.cell(1, 0).text = 'Possible values'
            table.cell(2, 0).text = 'Default value'
            table.cell(3, 0).text = 'Recommended value'

            i = 0
            while i < len(contexts):
                table.cell(i + 4, 0).text = contexts_columns[i]
                i+=1

            # # Check possible values
            if 'PossibleValues' in hardening_dataframe.columns:
                possible_values = policy['PossibleValues']
                if possible_values != '':
                    possible_values = possible_values.replace('[','').replace(']','').replace("'",'').split(",")
                    final_text = ''
                    i=0
                    for possible_value in possible_values:
                        if possible_value == 'nan':
                            final_text = ''
                        else:
                            if i != len(possible_values)-1:
                                final_text += '• ' + possible_value.strip() + '\n'
                            else:
                                final_text += '• ' + possible_value.strip()
                        i+=1
                else:
                    final_text = ''
            else:
                final_text = ''

            # Add cell data
            table.cell(1, 1).text = final_text
            table.cell(2, 1).text = policy['DefaultValue'] if policy['DefaultValue'] != 'nan' else ''
            table.cell(3, 1).text = policy['RecommendedValue'] if policy['RecommendedValue'] != 'nan' else ''

            i = 0
            while i < len(contexts):
                table.cell(i + 4, 1).text = policy[contexts[i]] if policy[contexts[i]] != 'nan' else ''
                i+=1
            
            # Vertical center cell text
            for cell in self.iter_cells(table):
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            # Place policy description
            if 'Description' in hardening_dataframe.columns and policy['Description'] != '':
                left = Inches(2)
                top = Inches(7.2)
                width = Inches(13)
                height = Inches(1.5)
                severity_rect = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
                severity_rect.shadow.inherit = False
                paragraph = severity_rect.text_frame.paragraphs[0]
                paragraph.alignment = PP_ALIGN.CENTER
                run = paragraph.add_run()
                run.text = policy['Description'].split("Note:",1)[0]
                run.font.name = "Calibri Light"
                run.font.size = Pt(16)
                run.font.color.rgb = RGBColor(0, 0, 0)
                severity_rect.fill.solid()
                lg_color = RGBColor(197, 224, 180) # light green

                severity_rect.fill.fore_color.rgb = lg_color
                severity_rect.line.color.rgb = lg_color

        prs.save(powerpoint_filepath)

    def create_applicable_csv(self, contexts: list, df_all_policies: pd.DataFrame):
        # retrieve policies and create csv
        try:
            CONTEXT_INDENT = 0
            for context in contexts:
                CONTEXT_INDENT += 1
                CONTEXT_NAME = "Context"+str(CONTEXT_INDENT)
                output_dataframe = pd.DataFrame(columns=df_all_policies.columns.values.tolist())
                for index, policy in context.iterrows():
                    # getting matching policy by name, .copy() 
                    # is to precise we want a copy and not a view of the dataframe
                    full_policy = df_all_policies.loc[df_all_policies['Name'] == policy['Name']].copy()
                    full_policy['RecommendedValue'] = policy[CONTEXT_NAME + ' - Computed Value']
                    output_dataframe = pd.concat([full_policy,output_dataframe.loc[:]]).reset_index(drop=True)
                    output_dataframe = output_dataframe.drop(['Description', 'Rationale', 'Impact'], axis=1)
                output_dataframe.to_csv(CONTEXT_NAME+".csv", index=False)
            return True
        except:
            return False

    def replace_defaults_values(self, output_csv):
        """
            This function will replace "DefaultValue" column content from
            file finding list to "-NODATA-"
        """
        df = pd.read_csv(self.file)
        df = df.fillna('')
        
        df =  df.assign(DefaultValue='-NODATA-')
        df.to_csv(output_csv, index=False)
        return True
    
    def get_number_of_context(self) -> int:
        """This function returns the number of contexts in an excel file

        Returns:
            int: number of contexts
        """
        df_contexts = pd.read_excel(self.file, "Contexts")
        context_number = 0
        for col in df_contexts.columns:
            if col.startswith('Context'):
                context_number+=1
        return context_number
    
    def get_contexts_names(self) -> list[str]:
        """This function return a list containing names of contexts in excel file

        Returns:
            list[str]: List of contexts names
        """
        df_contexts = pd.read_excel(self.file, "Contexts")
        return [col for col in df_contexts.columns if col.startswith('Context')]

    def create_xlsx(self) -> Workbook:
        """Create an excel file

        Returns:
            Workbook: The workbook object of the newly created excel file
        """
        workbook = Workbook()
        try:
            workbook.save(self.file)
        except:
            throw('An error occured while saving the Excel file, the name might be the cause.', 'high')
        workbook.load_workbook()
        return workbook
    
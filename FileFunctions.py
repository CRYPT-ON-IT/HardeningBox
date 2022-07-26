import os
import sys
import pandas as pd
import collections.abc
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.text import PP_ALIGN
from Errors import throw

class FileFunctions():

    def __init__(self, file):
        self.file = file

    """ 
        This function checks if a file exists and
        if the program can read the content.
    """
    def checkIfFileExistsAndReadable(self):
        if os.path.exists(self.file):
            print('\033[92mFile specified exists !\033[0m\n')
        else:
            throw('File specified not found, exiting.', 'high')

    """ 
        This function will return the content of
        a normal file.
    """
    def readFile(self):
        try:
            file = open(self.file, 'r')
            text = file.read()
            file.close()
        except:
            throw("Couldn't read file, exiting.", "high")
        
        return text

    """ 
        This function will return a dataframe (pandas)
        containing the whole data of a CSV file.
    """
    def readCsvFile(self):
        df = pd.read_csv(self.file, encoding='latin1')
        df = df.fillna('')
        df = df.astype(str)
        return df

    """ 
        This function will transform a CSV file
        into an Excel file, using pandas.
    """
    def convertCsv2Excel(self):
        df = pd.read_csv(self.file)
        df = df.fillna('')
        output_excel = ''
        output_excel_args = ['-o', '--output']
        for output_excel_arg in output_excel_args:
            for arg in sys.argv:
                if output_excel_arg == arg:
                    output_excel = sys.argv[sys.argv.index(arg)+1]
        if output_excel == '':
            output_excel = input("\nWhat's the name of the Excel output file ? : ")
        df.to_excel(output_excel, index=False)

    """ 
        This function will transform an Excel file
        into a CSV file, using pandas.
    """
    def convertExcel2Csv(self):
        df = pd.read_excel(self.file)
        df = df.fillna('')
        output_csv = ''
        output_csv_args = ['-o', '--output']
        for output_csv_arg in output_csv_args:
            for arg in sys.argv:
                if output_csv_arg == arg:
                    output_csv = sys.argv[sys.argv.index(arg)+1]
        if output_csv == '':
            output_csv = input("\nWhat's the name of the CSV output file ? : ")
        df.to_csv(output_csv, index=False)

    """ 
        This function will itterate over a PowerPoint
        table to transform the content.
    """
    def iter_cells(self, table):
        for row in table.rows:
            for cell in row.cells:
                yield cell
    
    """ 
        This function will transform a CSV file
        into PowerPoint Slides.
    """
    def CreatePPTX(self, hardening_dataframe, contexts, contexts_columns, powerpoint_filepath):

        # Creating presentation
        prs = Presentation()
        slide_size = (16, 9)
        prs.slide_width, prs.slide_height = Inches(slide_size[0]), Inches(slide_size[1])

        # Creating first slide 
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.width = Inches(16)
        title.height = Inches(1.5)
        title.top = Inches(3)
        subtitle.width = Inches(16)
        subtitle.height = Inches(1.5)
        subtitle.top = Inches(5)

        title.text = "Hardening presentation"
        subtitle.text = "Author"

        # Creating policies slides 
        for index, policy in hardening_dataframe.iterrows():
            tab_slide_layout = prs.slide_layouts[5]
            tab_slide = prs.slides.add_slide(tab_slide_layout)
            shapes = tab_slide.shapes

            shapes.title.text = policy['ID'] + ' ' + policy['Name']
            shapes.title.left = Inches(0)
            shapes.title.right = Inches(0)
            shapes.title.width = Inches(16)
            shapes.title.height = Inches(1.5)

            cols = 2
            rows = 5 + len(contexts)
            left = Inches(2.5)
            top = Inches(2.0)
            width = Inches(8.0)
            height = Inches(0.12)

            table = shapes.add_table(rows, cols, left, top, width, height).table

            # set column widths
            table.columns[1].width = Inches(8.8)

            # set rows height
            table.rows[0].height = Inches(0.7) # Title
            table.rows[1].height = Inches(0.7) # Possible values
            table.rows[2].height = Inches(0.7) # Default value
            table.rows[3].height = Inches(0.7) # Recommended value
            
            i = 0
            while i < len(contexts):
                table.rows[i + 4].height = Inches(0.7)
                i+=1

            table.rows[i + 4].height = Inches(0.7) # Comment

            # Default column headings
            table.cell(0, 0).text = 'Variable'
            table.cell(0, 1).text = 'Value'

            # Default 1st column cells
            table.cell(1, 0).text = 'Possible values'
            table.cell(2, 0).text = 'Default value'
            table.cell(3, 0).text = 'Recommended value'

            i = 0
            while i < len(contexts):
                table.cell(i + 4, 0).text = contexts_columns[i]
                i+=1

            table.cell(i + 4, 0).text = 'Comment'

            # Check possible values
            if 'PossibleValues' in hardening_dataframe.columns:
                possible_values = policy['PossibleValues']
                possible_values = possible_values.replace('[','').replace(']','').replace("'",'').split(",")
                final_text = ''
                i=0
                for possible_value in possible_values:
                    if possible_value == 'nan':
                        final_text = ''
                    else:
                        if i != len(possible_values)-1:
                            final_text += '• ' + possible_value.strip() + '\n'
                        else:
                            final_text += '• ' + possible_value.strip()
                    i+=1
            else:
                final_text = ''

            # Add cell data
            table.cell(1, 1).text = final_text
            table.cell(2, 1).text = policy['DefaultValue'] if policy['DefaultValue'] != 'nan' else ''
            table.cell(3, 1).text = policy['RecommendedValue'] if policy['RecommendedValue'] != 'nan' else ''

            i = 0
            while i < len(contexts):
                table.cell(i + 4, 1).text = policy[contexts[i]] if policy[contexts[i]] != 'nan' else ''
                i+=1

            table.cell(i + 4, 1).text = ''

            # Change font size
            for cell in self.iter_cells(table):
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(17)
            
            # Vertical center cell text
            for cell in self.iter_cells(table):
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            # Add Microsoft Link
            if 'MicrosoftLink' in hardening_dataframe.columns:
                left = Inches(1.0)
                top = Inches(8.2)
                width = Inches(15.0)
                height = Inches(0.5)
                text_box = shapes.add_textbox(left, top, width, height)
                paragraph = text_box.text_frame.paragraphs[0]
                paragraph.alignment = PP_ALIGN.CENTER
                run = paragraph.add_run()
                run.text = policy['MicrosoftLink'] if not policy['MicrosoftLink'] == 'nan' else ''
                run.hyperlink.address = policy['MicrosoftLink'] if not policy['MicrosoftLink'] == 'nan' else None
                text_box.text_frame.word_wrap = True
            

        prs.save(powerpoint_filepath)